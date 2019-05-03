from django.shortcuts import render
from django.urls import reverse
from django.contrib.auth.decorators import login_required
from django.http import HttpResponseRedirect,HttpResponse
from django.contrib.auth import authenticate,login,logout
from . import forms
from . import models
import multiprocessing
from pptx import Presentation
import sys
import os
import pandas as pd
import re
import zlib
import zipfile
        

# Create your views here.

## register user
def user_register(request):
    if request.method == 'POST':
            username     = request.POST.get('username')  
            password = request.POST.get('password')
            user_data = { 'username':username , 'password':password }
            user_form  = forms.UserForm(data=user_data)
            if user_form.is_valid():
                    user = user_form.save(commit=False)
                    user.set_password(user.password)
                    user.save()
                    return HttpResponseRedirect(reverse('homepage'))
            else:
                    form = forms.UserForm()
                    data = { 'form':form , 
                                'message':'Registration Failed, retry!!',
                                'status' : 'alert-danger' }
                    return render(request,'register.html',context=data)
    form = forms.UserForm()
    data = { 'form':form , 
                'message':'You are good to go.',
                'status' : 'alert-success' }
    return render(request,'register.html',context=data)




### login 
def user_login(request):
        if request.method == 'POST':
                username = request.POST.get('username')
                password = request.POST.get('password')
                user = authenticate(username=username,password=password)
                if user:
                        login(request,user)
                        return HttpResponseRedirect(reverse('homepage'))
                else:   
                        form = forms.UserForm()
                        data = { 'form':form , 
                                'message':'Email or Password wrong, retry!!',
                                'status' : 'alert-danger' }
                        return render(request,'login.html',context=data)
        else:   
                form = forms.UserForm()
                data = { 'form':form , 
                                'message':'You are good to go.',
                                'status' : 'alert-success' }
                return render(request,'login.html',context=data)


## logout
def user_logout(request):
	logout(request)
	return HttpResponseRedirect(reverse('homepage'))






### Generate Certificate Start
Path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
Path = Path + '/media/certificates/'

### Generate Pdf start 
def generatepdf(c_file):
        try:
                query = 'unoconv -f pdf '+ c_file
                os.system(query)
                os.remove(c_file)
        except:
            print('This Script Requires unoconv script on your computer')
            print("Couldn't Convert to pdf!!!")
## Generate Pdf End 


## generate Certificate start
def generate_cert(template,data,id):
        prs = Presentation(template)
        for slide in prs.slides:
                for shape in slide.shapes:
                        if not shape.has_text_frame:
                                continue
                        for paragraph in shape.text_frame.paragraphs:
                            for i in range(len(paragraph.runs)):
                                    label = paragraph.runs[i].text.lower()
                                    start = label.find('{')
                                    end = label.find('}')
                                    if start>=0 and end>0:
                                        label = label[start+1:end].strip()
                                        line = paragraph.runs[i].text
                                        paragraph.runs[i].text = line[:start] + data.get(label,'0') + line[end+1:]
                break

        prs.save(Path+str(id))
        generatepdf(Path+str(id))
## generate Certificate End


## certificate Main Process start
def process_project(project):
        # load xls start
        data_list = pd.read_excel(project.data,Sheetname='sheet1')
        labels = list(data_list.columns)
        total = len(data_list)
        project.total = total
        project.save()
        compression = zipfile.ZIP_DEFLATED
        zf = zipfile.ZipFile(Path+str(1169924 + project.project_id)+'.zip',mode='w')
        # load xls end
        #---------------------
        #generate certificate start
        for i in range(total):
                data={}
                for j in labels:
                    data.update({j:data_list[j.lower()][i]})
                id = str(1169924 + project.project_id) + "_" + str(project.current + 1)
                generate_cert(project.template,data,id)
                zf.write(Path+id+'.pdf',id+'.pdf',compress_type=compression)
                project.current = project.current + 1
                project.save()
        zf.close()
        project.link = str(1169924 + project.project_id)+'.zip'
        project.save()
        # generate certificate end
## certificate Main Process End
## Generate Certificate End


### verify Certificate start
def verify(certificate,projects):
        try:
                projectId = int(certificate.split('_')[0]) - 1169924
                certificateId = int(certificate.split('_')[1])
                for project in projects:
                        if project.project_id == projectId:
                                if certificateId > 0 and certificateId <= project.current:
                                        return { 'display':True,'verified' : True,'certificate':certificate}
                                else:
                                        pass
                return {'display':True,'verified':False, 'certificate':certificate}
        except:
                return {'display':True,'verified':False, 'certificate':certificate}
### verify Certificate End


### ger Certificate Start
def getCertificate(certificate,projects):
        try:
                projectId = int(certificate.split('_')[0]) - 1169924
                certificateId = int(certificate.split('_')[1])
                for project in projects:
                        if project.project_id == projectId:
                                if certificateId > 0 and certificateId <= project.current:
                                        link = certificate + '.pdf'
                                        return { 'display':True,'verified' : True,'certificate':certificate,'link':link}
                                else:
                                        pass
                return {'display':True,'verified':False, 'certificate':False}
        except:
                return {'display':True,'verified':False, 'certificate':certificate}
## get certificate End


### Homepage view
@login_required(login_url='/login')
def homepage(request):
        projects = list(models.ProjectData.objects.filter(user=request.user))
        verifiedCertificate = { 'display':False,'verified' : False,'certificate':False}     
        myCertificate = {'display': False}

        if request.method=='POST':
                form = forms.ProjectForm(request.POST,request.FILES)
                verifyCertificate = request.POST.get('verifyCertificate')
                certificateId = request.POST.get('certificateId')
                if verifyCertificate:
                        # certificate Verification
                        verifiedCertificate  = verify(verifyCertificate,projects)
                elif certificateId:
                        #get certficate
                        myCertificate = getCertificate(certificateId,projects)
                elif form.is_valid():
                        project = form.save(commit=False)
                        project.user = str(request.user)
                        project.save()
                        process = multiprocessing.Process(target=process_project,args=(project,))
                        process.start()
                        return HttpResponseRedirect(reverse('homepage'))
        
        project_form = forms.ProjectForm()
        completed = []
        ongoing = []
        for i in projects:
                if (i.total==0):
                        i.current = 0
                        ongoing.append(i)
                elif (i.current/i.total) == 1:
                        i.current = 100
                        completed.append(i)
                else:
                        i.current = int((i.current/i.total)*100)
                        ongoing.append(i)
        return render(request,'homepage.html',{'project_form_status':'','project_form':project_form,'username':request.user,'completed':completed,'ongoing':ongoing,'verified':verifiedCertificate,'certificate':myCertificate})

